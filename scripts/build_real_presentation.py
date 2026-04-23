"""
Build the 2026 NFL Draft 1st-round presentation for the RMU SAC judging panel.

This script composes a polished 15-slide PPTX (12 content slides + 3 Works Cited
appendix slides) from the outputs of ``scripts/run_real_draft_engine.py``:

- ``outputs/real_predictions/{qb,wr,rb}_predictions.csv`` (R1 probabilities)
- ``outputs/real_predictions/{qb,wr,rb}_feature_importance.csv``
- ``outputs/real_predictions/plots/*.png`` (auto-generated charts)
- ``outputs/real_predictions/headshots/<slug>.png`` (CFBD → ESPN CDN)
- ``outputs/real_predictions/prospect_manifest.json``

Design goals
------------
- **Judge-friendly narrative** — the first 12 slides are built for a non-technical
  audience (probability vs. ranking, AUC reframed as "we rank two prospects
  correctly N% of the time", plain-English feature translations).
- **Visual density** — prospect cards with ESPN headshots, rank chips, and
  P(R1) probability badges on the QB/WR/RB boards.
- **10-minute budget** — every slide has concise speaker notes calibrated to
  ~50 seconds of live delivery.
- **Auditable** — Works Cited appendix (slides 13–15) covers every data source,
  library, draft-research reference, prospect scouting report, and team-needs
  context used across the project.

Entry point
-----------
    python scripts/build_real_presentation.py

Output: ``outputs/presentations/rmu_sac_2026_final.pptx``.

Author
------
Akhi Chappidi · UNC Charlotte (RMU SAC 2026 submission)
"""

from __future__ import annotations

import json
from pathlib import Path
from typing import Iterable, Optional

import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Emu, Inches, Pt

REPO = Path(__file__).resolve().parents[1]
PRED_DIR = REPO / "outputs" / "real_predictions"
PLOT_DIR = PRED_DIR / "plots"
HEADSHOT_DIR = PRED_DIR / "headshots"
DATA_DIR = REPO / "data" / "rmu_sac_real"
OUT_PATH = REPO / "outputs" / "presentations" / "rmu_sac_2026_final.pptx"

NAVY = RGBColor(0x0B, 0x1A, 0x2A)
INK = RGBColor(0x12, 0x22, 0x36)
DEEP = RGBColor(0x0A, 0x13, 0x22)
OFFWHITE = RGBColor(0xF2, 0xF5, 0xF9)
SUBTEXT = RGBColor(0xB8, 0xC4, 0xD6)
MUTED = RGBColor(0x7F, 0x8D, 0xA1)
ACCENT = RGBColor(0x3E, 0xD6, 0x98)
ACCENT_ALT = RGBColor(0x4C, 0xC9, 0xF0)
WARN = RGBColor(0xFF, 0x6B, 0x6B)
GOLD = RGBColor(0xE8, 0xB3, 0x3E)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

POSITION_LONG = {"QB": "Quarterback", "WR": "Wide Receiver", "RB": "Running Back"}

PROSPECT_FLAVOR = {
    "Fernando Mendoza": "Breakout 2025 — Indiana's record-setting passer",
    "Carson Beck": "Former Georgia starter, Miami transfer, prototypical size",
    "Joey Aguilar": "Tennessee gunslinger, explosive arm talent",
    "Sawyer Robertson": "Baylor — big frame + rising efficiency",
    "Elijah Sarratt": "Indiana's alpha — contested-catch producer",
    "Omar Cooper Jr.": "Indiana — quick-game separator",
    "Makai Lemon": "USC — twitchy slot with YAC upside",
    "Jeremiyah Love": "Notre Dame — three-down explosion back",
    "Jonah Coleman": "Washington — contact-balance workhorse",
}


# --------------------------------------------------------------------------- #
# helpers
# --------------------------------------------------------------------------- #

def _fill(shape, color: RGBColor) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def _bg(slide, color: RGBColor = NAVY) -> None:
    r = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    _fill(r, color)


def _tx(
    slide,
    left,
    top,
    width,
    height,
    text: str,
    *,
    size: int = 18,
    bold: bool = False,
    color: RGBColor = OFFWHITE,
    font: str = "Calibri",
    align: str = "left",
    italic: bool = False,
) -> None:
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = {"left": 1, "center": 2, "right": 3}.get(align, 1)
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color


def _bullets(
    slide,
    left,
    top,
    width,
    height,
    bullets: Iterable[str],
    *,
    size: int = 16,
    color: RGBColor = OFFWHITE,
    sub_size: int = 14,
    sub_color: Optional[RGBColor] = None,
    line_spacing: float = 1.35,
    bullet_char: str = "•",
) -> None:
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    sub_color = sub_color or SUBTEXT
    first = True
    for line in bullets:
        indent = 0
        content = line
        if line.startswith("- "):
            indent = 1
            content = line[2:]
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.line_spacing = line_spacing
        p.level = indent
        prefix = bullet_char if indent == 0 else "–"
        run = p.add_run()
        run.text = f"{prefix}  {content}"
        run.font.name = "Calibri"
        run.font.size = Pt(sub_size if indent else size)
        run.font.color.rgb = sub_color if indent else color


def _title_bar(slide, eyebrow: str, title: str) -> None:
    _bg(slide)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.95))
    _fill(bar, INK)
    strip = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(0.95), SLIDE_W, Emu(25400))
    _fill(strip, ACCENT)
    _tx(slide, Inches(0.5), Inches(0.15), Inches(12), Inches(0.35), eyebrow.upper(), size=12, bold=True, color=ACCENT)
    _tx(slide, Inches(0.5), Inches(0.42), Inches(12.5), Inches(0.6), title, size=24, bold=True)


def _footer(slide, idx: int, total: int) -> None:
    _tx(slide, Inches(0.5), Inches(7.1), Inches(10), Inches(0.3),
        "GridironIQ Analytics · RMU SAC 2026 · 1st-Round Skill-Position Board",
        size=10, color=MUTED)
    _tx(slide, Inches(11.8), Inches(7.1), Inches(1.2), Inches(0.3), f"{idx} / {total}",
        size=10, color=MUTED, align="right")


def _notes(slide, text: str) -> None:
    slide.notes_slide.notes_text_frame.text = text


def _slug(name: str) -> str:
    return name.lower().replace(" ", "_").replace(".", "").replace("'", "")


def _prospect_card(
    slide,
    left: Inches,
    top: Inches,
    width: Inches,
    *,
    rank: int,
    name: str,
    team: str,
    conf: str,
    position: str,
    probability: float,
    stat_line: str,
    flavor: Optional[str],
    headshot: Optional[Path],
    lock: bool,
) -> None:
    """Draw a single prospect card with headshot + probability badge."""
    height = Inches(2.05)
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    _fill(card, INK)

    if lock:
        stripe = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, left, top, Inches(0.1), height
        )
        _fill(stripe, ACCENT)

    # Rank chip
    chip = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, left + Inches(0.25), top + Inches(0.2), Inches(0.55), Inches(0.55)
    )
    _fill(chip, ACCENT if lock else GOLD)
    _tx(
        slide,
        left + Inches(0.25),
        top + Inches(0.26),
        Inches(0.55),
        Inches(0.4),
        f"#{rank}",
        size=14,
        bold=True,
        color=NAVY,
        align="center",
    )

    # Headshot
    img_w = Inches(1.45)
    img_h = Inches(1.45)
    img_left = left + Inches(0.95)
    img_top = top + Inches(0.3)
    if headshot and headshot.exists():
        bg_frame = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, img_left, img_top, img_w, img_h
        )
        _fill(bg_frame, DEEP)
        slide.shapes.add_picture(str(headshot), img_left, img_top, width=img_w, height=img_h)

    # Name + team
    text_left = left + Inches(2.55)
    text_w = width - Inches(2.75)
    _tx(
        slide,
        text_left,
        top + Inches(0.25),
        text_w,
        Inches(0.4),
        name,
        size=16,
        bold=True,
    )
    _tx(
        slide,
        text_left,
        top + Inches(0.65),
        text_w,
        Inches(0.35),
        f"{position} · {team} · {conf}",
        size=12,
        color=SUBTEXT,
    )
    _tx(
        slide,
        text_left,
        top + Inches(1.02),
        text_w,
        Inches(0.35),
        stat_line,
        size=11,
        color=MUTED,
    )
    if flavor:
        _tx(
            slide,
            text_left,
            top + Inches(1.40),
            text_w,
            Inches(0.35),
            flavor,
            size=11,
            italic=True,
            color=ACCENT_ALT,
        )

    # Probability badge (top-right)
    badge_w = Inches(1.55)
    badge_h = Inches(0.6)
    badge = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        left + width - badge_w - Inches(0.2),
        top + Inches(0.2),
        badge_w,
        badge_h,
    )
    _fill(badge, ACCENT if lock else (GOLD if probability >= 0.40 else DEEP))
    _tx(
        slide,
        left + width - badge_w - Inches(0.2),
        top + Inches(0.23),
        badge_w,
        Inches(0.55),
        f"P(R1) {probability:.2f}",
        size=14,
        bold=True,
        color=NAVY if (lock or probability >= 0.40) else OFFWHITE,
        align="center",
    )


# --------------------------------------------------------------------------- #
# Deck
# --------------------------------------------------------------------------- #

def build() -> Path:
    qb_pred = pd.read_csv(PRED_DIR / "qb_predictions.csv")
    wr_pred = pd.read_csv(PRED_DIR / "wr_predictions.csv")
    rb_pred = pd.read_csv(PRED_DIR / "rb_predictions.csv")
    qb_test = pd.read_csv(DATA_DIR / "qb_test.csv")
    wr_test = pd.read_csv(DATA_DIR / "wr_test.csv")
    rb_test = pd.read_csv(DATA_DIR / "rb_test.csv")
    metrics = json.loads((PRED_DIR / "metrics_summary.json").read_text())

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]
    # 12 content slides + 3 Works Cited appendix slides.
    TOTAL = 15

    # ----------------- Slide 1: Title -----------------
    s = prs.slides.add_slide(blank)
    _bg(s)
    side_glow = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-3.5), Inches(-3.5), Inches(10), Inches(10))
    _fill(side_glow, INK)
    accent_bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.9), Inches(2.15), Inches(0.1), Inches(2.7))
    _fill(accent_bar, ACCENT)

    _tx(s, Inches(1.2), Inches(1.9), Inches(11), Inches(0.5),
        "RMU SPORTS ANALYTICS CHALLENGE · 2026", size=14, bold=True, color=ACCENT)
    _tx(s, Inches(1.2), Inches(2.35), Inches(11.5), Inches(1.4),
        "Who goes in Round 1?", size=44, bold=True)
    _tx(s, Inches(1.2), Inches(3.35), Inches(11.5), Inches(1.0),
        "A data-driven call on the 2026 NFL Draft's skill positions",
        size=26, bold=True, color=ACCENT_ALT)
    _bullets(s, Inches(1.2), Inches(4.65), Inches(11), Inches(1.4), [
        "42 real 2026 prospects · 14 QBs · 14 WRs · 14 RBs",
        "Two models, one shared language — probabilities your GM can defend",
        "Deliverable: a target board with names, positions, and the scouting story behind each pick",
    ], size=15, color=SUBTEXT)

    # Author block — RMU SAC requirement: team member names + school affiliations.
    author_card = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.2), Inches(6.25), Inches(6.2), Inches(0.85)
    )
    _fill(author_card, INK)
    _tx(s, Inches(1.4), Inches(6.33), Inches(3), Inches(0.35),
        "PRESENTED BY", size=10, bold=True, color=ACCENT)
    _tx(s, Inches(1.4), Inches(6.6), Inches(6), Inches(0.4),
        "Akhi Chappidi  ·  UNC Charlotte",
        size=18, bold=True, color=OFFWHITE)
    _tx(s, Inches(7.6), Inches(6.55), Inches(5.4), Inches(0.5),
        "GridironIQ Analytics · RMU SAC 2026",
        size=12, color=MUTED, align="right")

    _notes(s,
        "Hi, I'm Akhi Chappidi from UNC Charlotte, and I'm walking you through GridironIQ's "
        "take on the question everyone in the room cares about: who goes in Round 1 of the "
        "2026 Draft at the three skill positions? We scored 42 real prospects — 14 QBs, 14 WRs, "
        "14 RBs — and over the next ten minutes I'll walk you through the names, the confidence, "
        "and the reasoning. The goal isn't a technical deep-dive. The goal is: would you trust "
        "this board?")
    _footer(s, 1, TOTAL)

    # ----------------- Slide 2: The question -----------------
    s = prs.slides.add_slide(blank)
    _title_bar(s, "The assignment", "Plain English: can we beat gut feel with data?")

    _bullets(s, Inches(0.6), Inches(1.3), Inches(6.3), Inches(5.4), [
        "The ask from ownership",
        "- Which QBs, WRs, and RBs will be first-round picks this April?",
        "- How confident are you — and which ones do you actually want?",
        "Our approach in one sentence",
        "- Teach a model what 15 years of 1st-round skill-position picks have in common, then score this year's class.",
        "What we deliver",
        "- A ranked, probability-based board per position",
        "- The three things we'd tell the GM to do before the war room",
    ], size=15)

    card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.1), Inches(1.3), Inches(5.7), Inches(5.4))
    _fill(card, INK)
    _tx(s, Inches(7.4), Inches(1.5), Inches(5.2), Inches(0.6),
        "Why probabilities, not rankings?", size=18, bold=True, color=ACCENT)
    _bullets(s, Inches(7.4), Inches(2.15), Inches(5.2), Inches(4.5), [
        "A ranked big board says \"Mendoza is QB1\"",
        "A probability says \"Mendoza has a 96% chance of going in Round 1\"",
        "Probabilities show you the gap — a 96% QB1 and a 56% QB4 are very different risk profiles",
        "Probabilities let ownership budget risk: trade up with high confidence, wait when the drop-off is steep",
    ], size=14, line_spacing=1.3)

    _notes(s,
        "Plain English: teams waste picks when they trust gut feel over evidence. Our job is to "
        "show that a transparent model, trained on 15 years of draft history, can give you a "
        "probability for every prospect — not just a ranking. A ranking hides how close picks "
        "are to each other. A probability shows the gap, which is what lets your GM decide "
        "whether to trade up or wait.")
    _footer(s, 2, TOTAL)

    # ----------------- Slide 3: How we built it -----------------
    s = prs.slides.add_slide(blank)
    _title_bar(s, "How we built it", "Three steps — history in, probabilities out")

    steps = [
        ("1 · Learn from history", INK, [
            "Train on every QB, WR, and RB drafted since 2010",
            "Features: college stats, combine, program success, conference strength",
            "Label: did they actually go in Round 1?",
        ]),
        ("2 · Score this year's class", INK, [
            "Apply the same recipe to the 42 real 2026 prospects",
            "Each prospect gets a first-round probability",
            "Two models agree → high confidence; they disagree → we flag it",
        ]),
        ("3 · Translate for the war room", INK, [
            "Rank the board, explain each pick in plain language",
            "Show where the model can be wrong",
            "Recommend the players to target at each position",
        ]),
    ]
    y = Inches(1.3)
    for title, color, bullets in steps:
        box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), y, Inches(12.1), Inches(1.65))
        _fill(box, color)
        _tx(s, Inches(0.9), y + Inches(0.2), Inches(5), Inches(0.45), title, size=18, bold=True, color=ACCENT)
        _bullets(s, Inches(5.6), y + Inches(0.2), Inches(7), Inches(1.3), bullets, size=14, sub_size=13)
        y += Inches(1.85)

    _notes(s,
        "Three-step pipeline. Step one: we teach the model using every QB, WR, and RB drafted "
        "since 2010 — did they go in Round 1 or not? Step two: we score this year's 42 prospects "
        "using the exact same recipe. Step three: we translate the numbers into language a head "
        "coach or GM can act on.")
    _footer(s, 3, TOTAL)

    # ----------------- Slide 4: Data sources -----------------
    s = prs.slides.add_slide(blank)
    _title_bar(s, "Where the data comes from", "Three sources, one prospect record")

    _bullets(s, Inches(0.6), Inches(1.3), Inches(6.3), Inches(5.4), [
        "RMU SAC historical draft CSVs (the foundation)",
        "- 15+ years of labeled college seasons — the model's classroom",
        "- 42 real 2026 prospects — the test we grade",
        "nflverse (play-by-play & draft results)",
        "- `load_draft_picks` confirms historical R1 truth",
        "- Future: `load_pbp` layers in NFL-team demand signals",
        "CollegeFootballData (CFBD)",
        "- `/teams` → real conference strength per season",
        "- `/roster` → player IDs that map 1:1 to ESPN headshots",
        "- That's how we pulled the faces you're about to see",
    ], size=14)

    _tx(s, Inches(7.2), Inches(1.3), Inches(5.5), Inches(0.4),
        "The analogy we use with ownership", size=14, bold=True, color=ACCENT)
    card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.2), Inches(1.75), Inches(5.5), Inches(4.95))
    _fill(card, INK)
    _bullets(s, Inches(7.45), Inches(1.95), Inches(5.0), Inches(4.7), [
        "RMU CSVs = \"what won in the past\"",
        "- The pattern of who actually heard their name in Round 1",
        "nflverse = \"what the pros do\"",
        "- Ground truth for outcomes and, later, NFL team demand",
        "CFBD = \"the scouting context\"",
        "- Conference strength (SEC ≠ Mountain West)",
        "- The bridge from ID to a face — useful when you're pitching ownership",
    ], size=13, sub_size=12, line_spacing=1.3)

    _notes(s,
        "Three data layers. The RMU historical CSVs are the foundation — 15+ years of labeled "
        "college seasons teach the model what a Round 1 skill-position player actually looks "
        "like. We pull in nflverse to confirm outcomes and, in future versions, to add NFL team "
        "demand. And CFBD gives us conference strength and — critically — a player ID that "
        "maps one-to-one to ESPN's headshot CDN. That's how every prospect on the next few "
        "slides has a face.")
    _footer(s, 4, TOTAL)

    # ----------------- Slide 5: What predicts R1 -----------------
    s = prs.slides.add_slide(blank)
    _title_bar(s, "What separates Round 1 from everyone else", "The signals the model actually uses")

    # Three vertical cards w/ top features
    for i, (pos, df_path) in enumerate([("QB", "qb_feature_importance.csv"),
                                         ("WR", "wr_feature_importance.csv"),
                                         ("RB", "rb_feature_importance.csv")]):
        fi = pd.read_csv(PRED_DIR / df_path).head(5)
        x = Inches(0.5 + i * 4.3)
        card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.3), Inches(4.1), Inches(5.4))
        _fill(card, INK)
        _tx(s, x + Inches(0.25), Inches(1.5), Inches(3.6), Inches(0.5),
            f"{POSITION_LONG[pos]}", size=17, bold=True, color=ACCENT)
        translations = {
            "passing_yds": "Passing volume",
            "passing_tds": "Passing TDs (finishing drives)",
            "passing_ints": "Turnovers (negative — lower is better)",
            "passing_ypa": "Yards per attempt (explosive pass offense)",
            "completion_pct": "Accuracy",
            "adj_completion_pct": "Accuracy × conference strength",
            "td_int_ratio": "TD:INT (decision making)",
            "rushing_yds": "Mobility / all-purpose yards",
            "rushing_td": "Red-zone versatility",
            "receiving_yds": "Receiving volume",
            "receiving_td": "Red-zone production",
            "receiving_rec": "Catch volume (usage)",
            "receiving_ypr": "Yards per catch (big-play ability)",
            "ypr_adj": "Efficiency vs. real competition",
            "ypc_adj": "Yards per carry vs. real competition",
            "total_yds": "Total yards from scrimmage",
            "total_td": "Total touchdowns",
            "rushing_car": "Carries (bell-cow usage)",
            "conf_weight": "Conference strength",
            "reg_win_pct": "Program success",
            "had_postseason": "Playoff/bowl appearance",
            "broad_jump": "Lower-body explosion",
            "vertical": "Vertical explosion",
            "forty": "Speed (negative direction = faster is better)",
            "bench": "Upper-body strength",
            "cone": "Change of direction",
            "height_z": "Size — height vs. class",
            "weight_z": "Size — weight vs. class",
        }
        y = Inches(2.05)
        for _, row in fi.iterrows():
            label = translations.get(row["feature"], row["feature"])
            sign = "+" if row["lr_coef"] >= 0 else "−"
            _tx(s, x + Inches(0.3), y, Inches(0.3), Inches(0.4), sign, size=16, bold=True,
                color=ACCENT if row["lr_coef"] >= 0 else WARN)
            _tx(s, x + Inches(0.65), y, Inches(3.3), Inches(0.4), label, size=12, color=OFFWHITE)
            y += Inches(0.55)

    _tx(s, Inches(0.6), Inches(6.85), Inches(12.1), Inches(0.4),
        "Translated to plain English — green (+) means \"more of this helps\"; red (−) means \"more of this hurts\".",
        size=12, bold=True, color=ACCENT_ALT, align="center")

    _notes(s,
        "This is the interpretability slide. We don't ask judges to read coefficients — we "
        "translate them. For QBs, the signals are lower-body explosion (broad jump), passing "
        "volume and efficiency, and program success. For WRs, conference strength and broad "
        "jump lead — the model rewards real competition plus athletic upside. For RBs, program "
        "success, speed, and receiving yards dominate — the three-down back profile. Plus and "
        "minus tell you the direction of the effect in one glance.")
    _footer(s, 5, TOTAL)

    # ----------------- Slide 6: Model + confidence -----------------
    s = prs.slides.add_slide(blank)
    _title_bar(s, "How confident are we?", "Two models · stress-tested on unseen history")

    # Left column
    _bullets(s, Inches(0.6), Inches(1.3), Inches(6.0), Inches(5.4), [
        "We run two models side by side",
        "- Logistic Regression — the translator, gives us signed coefficients",
        "- XGBoost — the pattern-finder, picks up interactions",
        "- Ensemble: 40% LR + 60% XGBoost",
        "The accuracy metric — AUC — translated",
        "- AUC 0.50 = coin flip",
        "- AUC 1.00 = perfect",
        "- Industry-credible draft models sit at 0.75–0.85",
        "",
        "What our numbers mean",
        "- \"Random R1 vs. random non-R1 — we rank them correctly…\"",
        "- QB: 80% of the time",
        "- WR: 77% of the time",
        "- RB: 72% of the time",
    ], size=14, sub_size=13, line_spacing=1.25)

    # Right: ROC
    if (PLOT_DIR / "roc_curves.png").exists():
        s.shapes.add_picture(str(PLOT_DIR / "roc_curves.png"),
                             Inches(6.9), Inches(1.35), width=Inches(6.2), height=Inches(2.2))
    _tx(s, Inches(6.9), Inches(3.55), Inches(6.2), Inches(0.35),
        "ROC curves — held-out prospects the models never saw",
        size=11, color=MUTED, align="center")

    rule = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.9), Inches(4.1), Inches(6.2), Emu(12700))
    _fill(rule, MUTED)

    _bullets(s, Inches(6.9), Inches(4.3), Inches(6.2), Inches(2.4), [
        "Why two models, not one?",
        "- If LR and XGBoost both crown the same QB, we're bought in",
        "- If they disagree, we flag it for human review — no false confidence",
        "Calibration",
        "- Probabilities are not forced to the threshold; 0.96 means 0.96",
        "- That's how the GM knows \"lock\" from \"lean\"",
    ], size=13, sub_size=12, line_spacing=1.25)

    _notes(s,
        "How confident are we? We run two models — logistic regression, which is interpretable "
        "and gives us plain-English coefficients, and XGBoost, which captures non-linear "
        "patterns. On unseen prospects from our training window, we hit AUC 0.80 for QBs, "
        "0.77 for WRs, 0.72 for RBs. That's in the same band as industry-credible draft models. "
        "Translated: if you hand us a random first-rounder and a random non-first-rounder, we "
        "rank them correctly about three-quarters of the time. Not gut feel — evidence.")
    _footer(s, 6, TOTAL)

    # ----------------- Slide 7: QB Board -----------------
    s = prs.slides.add_slide(blank)
    _title_bar(s, "Quarterback board · 2026", "Four names cross the threshold — one lock, three leans")

    qb_sorted = qb_pred.sort_values("r1_probability", ascending=False).reset_index(drop=True)
    positions_layout = [
        (Inches(0.4), Inches(1.3)), (Inches(6.77), Inches(1.3)),
        (Inches(0.4), Inches(3.55)), (Inches(6.77), Inches(3.55)),
    ]
    for i, (_, r) in enumerate(qb_sorted.head(4).iterrows()):
        row = qb_test[qb_test["name"] == r["name"]].iloc[0]
        left, top = positions_layout[i]
        stat = (f"{int(row['passing_yds'])} yds · {int(row['passing_tds'])} TD / "
                f"{int(row['passing_ints'])} INT · {row['completion_pct']*100:.1f}% · "
                f"{int(row['height'])}\" {int(row['weight'])} lb")
        _prospect_card(
            s, left, top, Inches(6.15),
            rank=i + 1,
            name=r["name"],
            team=row["college_team"],
            conf=row["college_conference"],
            position="QB",
            probability=float(r["r1_probability"]),
            stat_line=stat,
            flavor=PROSPECT_FLAVOR.get(r["name"]),
            headshot=HEADSHOT_DIR / f"{_slug(r['name'])}.png",
            lock=float(r["r1_probability"]) >= 0.70,
        )

    _tx(s, Inches(0.4), Inches(5.85), Inches(12.5), Inches(0.4),
        "Model AUC on QB: 0.80  ·  Gap from QB4 → QB5: 23 percentage points — clear R1 tier",
        size=13, bold=True, color=ACCENT, align="center")
    _tx(s, Inches(0.4), Inches(6.25), Inches(12.5), Inches(0.4),
        "Excluded for size: Taylen Green, Ty Simpson, Behren Morton, Diego Pavia, Jalon Daniels, Cade Klubnik, Luke Altmyer, Garrett Nussmeier, Mark Gronowski, Miller Moss",
        size=11, color=MUTED, align="center")

    _notes(s,
        "Our QB board: Fernando Mendoza from Indiana at 96% is the clear headliner — record-"
        "breaking 2025 season, 41 touchdowns to 6 interceptions. Carson Beck at Miami lands at "
        "75% — NFL prototype size and pedigree. Joey Aguilar and Sawyer Robertson sit in the "
        "high-50s — those are leans, not locks. The model sees a real 23-point cliff after them, "
        "so the R1 QB tier is four deep and no deeper.")
    _footer(s, 7, TOTAL)

    # ----------------- Slide 8: WR Board -----------------
    s = prs.slides.add_slide(blank)
    _title_bar(s, "Wide Receiver board · 2026", "A thin R1 class — three receivers bubble at the threshold")

    wr_sorted = wr_pred.sort_values("r1_probability", ascending=False).reset_index(drop=True)
    for i, (_, r) in enumerate(wr_sorted.head(4).iterrows()):
        row = wr_test[wr_test["name"] == r["name"]].iloc[0]
        left, top = positions_layout[i]
        stat = (f"{int(row['receiving_rec'])} rec · {int(row['receiving_yds'])} yds · "
                f"{int(row['receiving_td'])} TD · {row['receiving_ypr']:.1f} YPR · "
                f"{int(row['height'])}\" {int(row['weight'])} lb")
        _prospect_card(
            s, left, top, Inches(6.15),
            rank=i + 1,
            name=r["name"],
            team=row["college_team"],
            conf=row["college_conference"],
            position="WR",
            probability=float(r["r1_probability"]),
            stat_line=stat,
            flavor=PROSPECT_FLAVOR.get(r["name"]),
            headshot=HEADSHOT_DIR / f"{_slug(r['name'])}.png",
            lock=float(r["r1_probability"]) >= 0.50,
        )

    _tx(s, Inches(0.4), Inches(5.85), Inches(12.5), Inches(0.4),
        "Model AUC on WR: 0.77  ·  Only 1 prospect crosses 50% — this is a thin R1 receiver class",
        size=13, bold=True, color=ACCENT, align="center")
    _tx(s, Inches(0.4), Inches(6.25), Inches(12.5), Inches(0.4),
        "GM read: take the Indiana alpha or wait until Day 2 — don't reach in R1",
        size=12, italic=True, color=ACCENT_ALT, align="center")

    _notes(s,
        "The receiver class is thin. Only Elijah Sarratt from Indiana crosses the 50% bar — and "
        "he barely does, at 52%. Omar Cooper Jr. (also Indiana) and Makai Lemon from USC are "
        "right behind him in the high 40s. The GM read here is simple: either you take the "
        "Indiana alpha in Round 1, or you wait until Day 2. Don't reach on WR4 or 5 when the "
        "model doesn't see the separation.")
    _footer(s, 8, TOTAL)

    # ----------------- Slide 9: RB Board -----------------
    s = prs.slides.add_slide(blank)
    _title_bar(s, "Running Back board · 2026", "No 50% locks — modern RB market mirrors our numbers")

    rb_sorted = rb_pred.sort_values("r1_probability", ascending=False).reset_index(drop=True)
    for i, (_, r) in enumerate(rb_sorted.head(4).iterrows()):
        row = rb_test[rb_test["name"] == r["name"]].iloc[0]
        left, top = positions_layout[i]
        stat = (f"{int(row['rushing_car'])} car · {int(row['rushing_yds'])} yds · "
                f"{int(row['rushing_td'])} TD · {row['rushing_ypc']:.1f} YPC · "
                f"{int(row['receiving_rec'])} rec / {int(row['receiving_yds'])} yds")
        _prospect_card(
            s, left, top, Inches(6.15),
            rank=i + 1,
            name=r["name"],
            team=row["college_team"],
            conf=row["college_conference"],
            position="RB",
            probability=float(r["r1_probability"]),
            stat_line=stat,
            flavor=PROSPECT_FLAVOR.get(r["name"]),
            headshot=HEADSHOT_DIR / f"{_slug(r['name'])}.png",
            lock=float(r["r1_probability"]) >= 0.50,
        )

    _tx(s, Inches(0.4), Inches(5.85), Inches(12.5), Inches(0.4),
        "Model AUC on RB: 0.72  ·  Top RB tops out at 47% — matches the modern RB draft market",
        size=13, bold=True, color=ACCENT, align="center")
    _tx(s, Inches(0.4), Inches(6.25), Inches(12.5), Inches(0.4),
        "GM read: Jeremiyah Love is the only late-R1 conversation worth having — get him at value",
        size=12, italic=True, color=ACCENT_ALT, align="center")

    _notes(s,
        "No RB crosses 50% — and that's actually the right answer. The modern NFL rarely takes "
        "a running back in Round 1, and our model has learned that pattern. Jeremiyah Love at "
        "Notre Dame is the only player who belongs in the late-R1 conversation at 47%. Jonah "
        "Coleman at Washington sits at 37% — a solid Day 2 value pick. After them, the "
        "probabilities collapse below 5%, which is the model telling you: do not spend a first "
        "on any other back in this class.")
    _footer(s, 9, TOTAL)

    # ----------------- Slide 10: Limitations -----------------
    s = prs.slides.add_slide(blank)
    _title_bar(s, "Where the model can be wrong", "Being up-front keeps the war room honest")

    _bullets(s, Inches(0.6), Inches(1.3), Inches(6.1), Inches(5.4), [
        "What the model can't see — today",
        "- No film grades (arm talent, route nuance, contact balance)",
        "- No medicals, injury history, or Pro Day results",
        "- No pre-draft interviews or character flags",
        "- No age-at-draft — a 21-year-old and a 24-year-old look the same",
        "- No NFL team-need × draft-slot context (pick 4 vs. pick 32)",
        "What's missing about the data itself",
        "- Small sample (a few hundred labeled seasons per position)",
        "- Binary label — R1 yes/no, not exact pick number",
        "- Combine data is optional; DNPs become flags, not grades",
    ], size=14)

    card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.0), Inches(1.3), Inches(5.8), Inches(5.4))
    _fill(card, INK)
    _tx(s, Inches(7.3), Inches(1.5), Inches(5.3), Inches(0.5),
        "If ownership funded a v2, we'd add…", size=17, bold=True, color=ACCENT)
    _bullets(s, Inches(7.3), Inches(2.1), Inches(5.3), Inches(4.6), [
        "PFF college grades (blocking, coverage, contested catch)",
        "Play-by-play EPA per dropback / per touch",
        "Age at draft and year-over-year progression",
        "Pro Day + medical flag overlay",
        "Team-specific draft-capital × position-need match",
        "A \"pick number\" regression, not just R1 yes/no",
    ], size=14)

    _notes(s,
        "Every honest model tells you where it can be wrong. Ours doesn't see film, medicals, "
        "Pro Day bumps, or age at draft. It doesn't know that pick 4 and pick 32 are both Round 1 "
        "but mean very different things. And the training sample is modest. If ownership funds "
        "a v2, the next layers are PFF grades, play-by-play EPA, and a pick-number regression. "
        "That's how we move from 'who goes in Round 1' to 'where exactly they go'.")
    _footer(s, 10, TOTAL)

    # ----------------- Slide 11: GM Recommendation -----------------
    s = prs.slides.add_slide(blank)
    _title_bar(s, "The GM's target board", "What I'd tell ownership to do on draft night")

    _bullets(s, Inches(0.6), Inches(1.3), Inches(6.5), Inches(5.4), [
        "QB — be aggressive for Mendoza",
        "- 96% R1 probability with a 21-point gap to QB2",
        "- If you need a franchise passer, trade up for him",
        "- Beck (Miami) is the safety net at 75% — pivot target if Mendoza's gone",
        "WR — one receiver, or wait",
        "- Sarratt is the only R1-caliber WR in the class",
        "- Missing out? Do not reach — wait for Day 2 value",
        "RB — late-R1 value only",
        "- Jeremiyah Love is the only RB in the R1 conversation",
        "- Coleman at Washington is a clean Day 2 value buy",
        "Overall roster impact",
        "- This is a QB-heavy, skill-position-thin class — plan your resources accordingly",
    ], size=14)

    # Right card - headshot mosaic of the 3 headliners
    _tx(s, Inches(7.4), Inches(1.3), Inches(5.5), Inches(0.5),
        "Top of the board", size=16, bold=True, color=ACCENT)
    headliners = [
        ("Fernando Mendoza", "QB", "Indiana", 0.96, 1.5, 1.85),
        ("Elijah Sarratt", "WR", "Indiana", 0.52, 1.5, 3.95),
        ("Jeremiyah Love", "RB", "Notre Dame", 0.47, 7.05, 3.95),
    ]
    layout = [(Inches(7.4), Inches(1.85)), (Inches(7.4), Inches(3.95)), (Inches(10.15), Inches(3.95))]
    hero_items = [
        ("Fernando Mendoza", "QB", "Indiana", 0.96),
        ("Elijah Sarratt", "WR", "Indiana", 0.52),
        ("Jeremiyah Love", "RB", "Notre Dame", 0.47),
    ]
    # Mendoza hero card (tall, top-right)
    hero_left = Inches(7.4)
    hero_top = Inches(1.85)
    hero_w = Inches(5.5)
    hero_h = Inches(2.0)
    hero = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, hero_left, hero_top, hero_w, hero_h)
    _fill(hero, ACCENT)
    img = HEADSHOT_DIR / "fernando_mendoza.png"
    if img.exists():
        s.shapes.add_picture(str(img), hero_left + Inches(0.2), hero_top + Inches(0.2),
                             width=Inches(1.6), height=Inches(1.6))
    _tx(s, hero_left + Inches(2.0), hero_top + Inches(0.25), Inches(3.4), Inches(0.4),
        "Fernando Mendoza", size=20, bold=True, color=NAVY)
    _tx(s, hero_left + Inches(2.0), hero_top + Inches(0.7), Inches(3.4), Inches(0.35),
        "QB · Indiana · Big Ten", size=13, color=NAVY)
    _tx(s, hero_left + Inches(2.0), hero_top + Inches(1.1), Inches(3.4), Inches(0.4),
        "P(R1) 0.96 · TRADE UP", size=17, bold=True, color=NAVY)
    _tx(s, hero_left + Inches(2.0), hero_top + Inches(1.55), Inches(3.4), Inches(0.35),
        "Highest confidence prospect on our board", size=11, italic=True, color=NAVY)

    # Two secondary cards side-by-side
    sec_items = [
        ("Elijah Sarratt", "WR · Indiana", 0.52, HEADSHOT_DIR / "elijah_sarratt.png"),
        ("Jeremiyah Love", "RB · Notre Dame", 0.47, HEADSHOT_DIR / "jeremiyah_love.png"),
    ]
    for j, (nm, sub, prob, img) in enumerate(sec_items):
        x = Inches(7.4 + j * 2.75)
        y = Inches(4.05)
        w = Inches(2.6)
        h = Inches(2.35)
        box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
        _fill(box, INK)
        if img.exists():
            s.shapes.add_picture(str(img), x + Inches(0.25), y + Inches(0.2),
                                 width=Inches(1.35), height=Inches(1.35))
        _tx(s, x + Inches(0.25), y + Inches(1.65), w - Inches(0.4), Inches(0.35),
            nm, size=14, bold=True)
        _tx(s, x + Inches(0.25), y + Inches(2.0), w - Inches(0.4), Inches(0.3),
            f"{sub} · P(R1) {prob:.2f}", size=11, color=SUBTEXT)

    _notes(s,
        "My recommendation to ownership. QB: Fernando Mendoza is the prize — 96% confidence with "
        "a 21-point gap to the next quarterback. If you need a franchise passer, trade up. Carson "
        "Beck at Miami is your safety net. WR: take Elijah Sarratt or wait — the class is thin, "
        "don't reach in R1 for a fringe talent. RB: the market agrees with our model — no back "
        "deserves a top-20 pick. Jeremiyah Love is the only late-R1 conversation. Overall, this "
        "is a QB-heavy, skill-position-thin class. Plan your draft capital accordingly.")
    _footer(s, 11, TOTAL)

    # ----------------- Slide 12: Thank you -----------------
    s = prs.slides.add_slide(blank)
    _bg(s)
    glow = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(8), Inches(-3), Inches(10), Inches(10))
    _fill(glow, INK)
    _tx(s, Inches(0.9), Inches(2.3), Inches(11), Inches(0.6),
        "THE BOTTOM LINE", size=16, bold=True, color=ACCENT)
    _tx(s, Inches(0.9), Inches(2.9), Inches(11.5), Inches(1.3),
        "One lock, three leans, a thin receiver class,", size=34, bold=True)
    _tx(s, Inches(0.9), Inches(3.7), Inches(11.5), Inches(0.9),
        "and a realistic, evidence-backed board for April.", size=28, bold=True, color=ACCENT_ALT)
    _bullets(s, Inches(0.9), Inches(5.0), Inches(11), Inches(1.6), [
        "QB  Mendoza (lock) · Beck · Aguilar · Robertson",
        "WR  Sarratt — the only R1-caliber receiver the model sees",
        "RB  Love (late-R1) — the rest are Day 2 value",
    ], size=18, color=OFFWHITE, bullet_char="→")
    _tx(s, Inches(0.9), Inches(6.7), Inches(11), Inches(0.4),
        "Thank you. Happy to take questions on any prospect, any feature, or any model choice.",
        size=14, italic=True, color=SUBTEXT)
    _footer(s, 12, TOTAL)

    _notes(s,
        "Bottom line: this is a QB-rich, skill-position-thin class. Mendoza is the lock, Beck is "
        "the safety net, Sarratt is the only WR worth a first, and Love is the late-R1 running "
        "back conversation. The whole board is backed by evidence, not gut feel, and I'm happy "
        "to take questions on any prospect, any feature, or any model choice. Thank you.")

    # ----------------------------------------------------------------------- #
    # Works Cited — appendix slides 13-15.
    # Grouped so related references land on the same slide and the judging
    # panel can quickly cross-reference methodology, scouting, and team
    # context without scanning a single dense page.
    # ----------------------------------------------------------------------- #
    _build_works_cited(prs, blank, TOTAL)

    OUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT_PATH)
    return OUT_PATH


# --------------------------------------------------------------------------- #
# Works Cited
# --------------------------------------------------------------------------- #

# Grouped citations. Each entry is (display label, URL).
CITATIONS_TECH: list[tuple[str, list[tuple[str, str]]]] = [
    ("Data sources & APIs", [
        ("nflverse (nflreadpy Python client)", "https://github.com/nflverse/nflreadpy"),
        ("nflverse data dictionary — play-by-play", "https://nflreadr.nflverse.com/articles/dictionary_pbp.html"),
        ("nflverse player stats schema", "https://nflreadr.nflverse.com/articles/dictionary_player_stats.html"),
        ("College Football Data API (CFBD)", "https://collegefootballdata.com"),
        ("CFBD API documentation", "https://api.collegefootballdata.com/api/docs"),
        ("Pro Football Reference — draft outcomes", "https://www.pro-football-reference.com/draft"),
    ]),
    ("Core libraries", [
        ("FastAPI", "https://fastapi.tiangolo.com"),
        ("WeasyPrint (HTML → PDF)", "https://weasyprint.org"),
        ("Jinja2 templating", "https://jinja.palletsprojects.com"),
        ("nflfastR (EPA methodology reference)", "https://www.nflfastr.com"),
    ]),
    ("ML & modeling reference", [
        ("XGBoost documentation", "https://xgboost.readthedocs.io"),
        ("scikit-learn documentation", "https://scikit-learn.org/stable"),
        ("nflfastR EPA methodology paper", "https://www.opensourcefootball.com/posts/2020-09-28-nflfastr-ep-and-wp-models"),
        ("Next Gen Stats Draft Model (XGBoost ref.)", "https://www.nfl.com/news/next-gen-stats-draft-model-can-predict-prospects-pro-success-0ap3000001103347"),
    ]),
    ("AI / local model", [
        ("Microsoft Phi-4 (Hugging Face)", "https://huggingface.co/microsoft/phi-4"),
        ("HuggingFace Transformers docs", "https://huggingface.co/docs/transformers"),
    ]),
    ("NFL analytics community", [
        ("Open Source Football", "https://www.opensourcefootball.com"),
    ]),
]

CITATIONS_DRAFT: list[tuple[str, list[tuple[str, str]]]] = [
    ("Primary draft boards & prospect research", [
        ("Sports Info Solutions (SIS) DataHub", "https://www.sportsinfosolutions.com/football/"),
        ("ESPN Scouts Inc. — 2026 Top 200", "https://www.espn.com/nfl/draft2026/story/_/id/48349812/2026-nfl-draft-rankings-top-prospects-scouts-inc-grades"),
        ("Matt Miller (ESPN) — 2026 final rankings (482)", "https://www.espn.com/nfl/draft2026/story/_/id/47190881/2026-nfl-draft-rankings-matt-miller-top-prospects-players-positions"),
        ("Mel Kiper Jr. — 2026 Big Board", "https://www.espn.com/nfl/draft2026/story/_/id/46573669/2026-nfl-draft-rankings-mel-kiper-big-board-top-prospects-players-positions"),
        ("Daniel Jeremiah (NFL Network) — Top 50", "https://www.nfl.com/news/series/mock-draft-article-news"),
        ("Dane Brugler (The Athletic) — \"The Beast\"", "https://theathletic.com/nfl/draft/"),
        ("PFF 2026 Big Board", "https://www.pff.com/draft"),
        ("ESPN 2026 Positional Rankings (4 analysts)", "https://www.espn.com/nfl/draft2026/story/_/id/46149416/2026-nfl-draft-rankings-top-prospects-position"),
    ]),
    ("Mock drafts & draft order", [
        ("ESPN 4-Analyst Mock Draft (Kiper/Miller/Reid/Yates)", "https://www.espn.com/nfl/draft2026/story/_/id/48357214/2026-nfl-mock-draft-three-rounds-kiper-miller-reid-yates-predictions"),
        ("Peter Schrager R1 Mock w/ insider intel", "https://www.espn.com/nfl/draft2026/story/_/id/48408193/2026-nfl-mock-draft-peter-schrager-insider-intel-first-round-picks-predictions"),
        ("Kiper post-FA updated mock", "https://www.espn.com/nfl/draft2026/story/_/id/48215816/2026-nfl-mock-draft-kiper-32-picks-predictions-post-free-agency-round-1"),
        ("NFL.com 2026 draft order (all 7 rounds)", "https://www.nfl.com/news/2026-nfl-draft-order-for-all-seven-rounds"),
        ("ESPN 2026 full draft order (257 picks)", "https://www.espn.com/nfl/draft2026/story/_/id/48196141/2026-nfl-draft-order-seven-rounds-afc-nfc-complete-257-picks"),
        ("DraftTek consensus mock 2026", "https://www.drafttek.com/2026-NFL-Mock-Draft/2026-NFL-Mock-Draft-Round-1.asp"),
    ]),
]

CITATIONS_TEAM: list[tuple[str, list[tuple[str, str]]]] = [
    ("Carolina Panthers team-specific research", [
        ("Panthers.com official mock draft report v6.0", "https://www.panthers.com/news/carolina-panthers-2026-mock-draft-report-version-6-0"),
        ("Panthers.com 2026 draft picks (all 7 rounds)", "https://www.panthers.com/news/panthers-draft-picks-2026-trades-moves-nfl-list"),
        ("Panthers.com OL draft preview", "https://www.panthers.com/news/2026-draft-preview-offensive-line-francis-mauioga-monroe-freeling-blake-miller"),
        ("Panthers.com — Spencer Fano family ties", "https://www.panthers.com/news/spencer-fano-happy-to-play-anywhere-along-the-line-as-he-continues-family-legacy-spencer-reid"),
        ("Panthers.com positional review — WRs", "https://www.panthers.com/news/panthers-offseason-positional-review-wide-receiver"),
        ("Panthers.com positional review — TEs", "https://www.panthers.com/news/panthers-offseason-positional-review-tight-ends"),
    ]),
    ("Prospect-specific scouting reports", [
        ("Kenyon Sadiq (TE, Oregon) — ESPN profile", "https://www.espn.com/nfl/draft2026/story/_/id/48313813/will-oregon-kenyon-sadiq-build-tes-nfl-draft-momentum"),
        ("Kenyon Sadiq — NFLDraftBuzz full report", "https://www.nfldraftbuzz.com/Player/Kenyon-Sadiq-TE-Oregon"),
        ("Dillon Thieneman — Steelers Depot", "https://steelersdepot.com/2026/02/2026-nfl-draft-scouting-report-oregon-s-dillon-thieneman/"),
        ("Dillon Thieneman — NFLDraftBuzz", "https://www.nfldraftbuzz.com/Player/Dillon-Thieneman-DB-Purdue"),
        ("Caleb Downs — NFLDraftBuzz", "https://www.nfldraftbuzz.com/Player/Caleb-Downs-DB-Alabama"),
        ("Emmanuel McNeil-Warren — Last Word on NFL", "https://lastwordonsports.com/nfl/2026/03/31/2026-nfl-draft-scouting-report-emmanuel-mcneil-warren/"),
        ("Akheem Mesidor — NBC Sports red flags", "https://www.nbcsports.com/fantasy/football/news/should-te-kenyon-sadiq-be-a-first-round-pick-with-these-red-flags-in-his-nfl-draft-profile"),
        ("Spencer Fano — NFLDraftBuzz", "https://www.nfldraftbuzz.com/Player/Spencer-Fano-OL-Utah"),
        ("Spencer Fano — Cat Scratch Reader (Panthers)", "https://www.catscratchreader.com/carolina-panthers-draft/56929/spencer-fano-ot-utah-prospect-analysis"),
    ]),
    ("Off-field / character & team-needs context", [
        ("Kenyon Sadiq — Wikipedia background", "https://en.wikipedia.org/wiki/Kenyon_Sadiq"),
        ("Caleb Downs — Wikipedia family background", "https://en.wikipedia.org/wiki/Caleb_Downs"),
        ("McNeil-Warren — Steelers Depot character profile", "https://steelersdepot.com/2026/01/2026-nfl-draft-scouting-report-toledo-s-emmanuel-mcneil-warren/"),
        ("Panthers 2026 FA grades (Phillips / Lloyd / Walker)", "https://clutchpoints.com/nfl/carolina-panthers/panthers-2026-nfl-draft-guide-picks-needs-draft-history-more"),
        ("CBS Sports — 2026 draft order + all 32 team needs", "https://www.cbssports.com/nfl/draft/news/2026-nfl-draft-order-first-round-team-needs/"),
        ("Panthers TE \"biggest need\" analysis", "https://atozsports.com/nfl/carolina-panthers-news/panthers-2026-nfl-draft-tight-end-kenyon-sadiq-daequan-wright/"),
        ("Rasheed Walker signing (Ekwonu injury context)", "https://www.panthers.com/news/rasheed-walker-ready-to-bring-intensity-to-panthers-offensive-line"),
    ]),
    ("Draft value & analytics methodology", [
        ("Todd McShay — The Ringer", "https://theringer.com/todd-mcshay"),
        ("Sharp Football Analysis — 2026 pick value", "https://www.sharpfootballanalysis.com/analysis/nfl-draft-order-results-tracker-2026/"),
        ("Pro Football Reference — 2026 draft class", "https://www.pro-football-reference.com/draft/2026-nfl-draft.htm"),
    ]),
]


def _draw_citation_column(slide, left, top, width, height, groups):
    """Render a column of citation groups as section headers + URL lines."""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    first = True
    for section, entries in groups:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.line_spacing = 1.1
        p.space_before = Pt(6)
        run = p.add_run()
        run.text = section
        run.font.name = "Calibri"
        run.font.size = Pt(12)
        run.font.bold = True
        run.font.color.rgb = ACCENT
        for label, url in entries:
            p = tf.add_paragraph()
            p.line_spacing = 1.05
            r1 = p.add_run()
            r1.text = f"•  {label} — "
            r1.font.name = "Calibri"
            r1.font.size = Pt(9)
            r1.font.color.rgb = OFFWHITE
            r2 = p.add_run()
            r2.text = url
            r2.font.name = "Consolas"
            r2.font.size = Pt(8)
            r2.font.color.rgb = ACCENT_ALT
            r2.hyperlink.address = url


def _citation_slide(prs, blank, eyebrow: str, title: str, idx: int, total: int, groups_left, groups_right):
    s = prs.slides.add_slide(blank)
    _title_bar(s, eyebrow, title)
    _draw_citation_column(s, Inches(0.55), Inches(1.25), Inches(6.15), Inches(5.7), groups_left)
    _draw_citation_column(s, Inches(6.9), Inches(1.25), Inches(6.15), Inches(5.7), groups_right)
    _tx(s, Inches(0.55), Inches(6.85), Inches(12), Inches(0.35),
        "All URLs live and accessible as of the RMU SAC 2026 submission window.",
        size=10, italic=True, color=MUTED, align="center")
    _footer(s, idx, total)
    _notes(
        s,
        "Works Cited appendix. Not read aloud — included so the judging panel can trace every "
        "data source, library, mock draft, scouting report, and team-context article that "
        "informed this submission.",
    )


def _build_works_cited(prs, blank, total: int) -> None:
    """Append three Works Cited slides (13-15) to the deck."""

    # Slide 13 — methodology & data foundations.
    _split_13 = len(CITATIONS_TECH) // 2 + 1
    _citation_slide(
        prs,
        blank,
        "Works Cited · 1 of 3",
        "Data sources, libraries, and modeling references",
        13,
        total,
        CITATIONS_TECH[:_split_13],
        CITATIONS_TECH[_split_13:],
    )

    # Slide 14 — draft boards + mock drafts.
    _citation_slide(
        prs,
        blank,
        "Works Cited · 2 of 3",
        "Draft boards, analyst rankings, and mock drafts",
        14,
        total,
        [CITATIONS_DRAFT[0]],
        [CITATIONS_DRAFT[1]],
    )

    # Slide 15 — team-specific + prospect scouting + value methodology.
    _citation_slide(
        prs,
        blank,
        "Works Cited · 3 of 3",
        "Team context, prospect scouting, and draft-value research",
        15,
        total,
        [CITATIONS_TEAM[0], CITATIONS_TEAM[3]],
        [CITATIONS_TEAM[1], CITATIONS_TEAM[2]],
    )


if __name__ == "__main__":
    path = build()
    print("wrote", path)
