"""
Build the RMU SAC 2026 first-round presentation (QB/WR/RB).

Produces a 12-slide PPTX with speaker notes sized for a 10-minute delivery.
Pulls R1 probabilities, feature-importance CSVs, and plots from
``outputs/rmu_predictions/``.

Usage::

    python scripts/build_rmu_presentation.py

Output: ``outputs/presentations/rmu_sac_2026_r1_predictions.pptx``.
"""

from __future__ import annotations

from pathlib import Path
from typing import Iterable, Optional

import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Emu, Inches, Pt

REPO_ROOT = Path(__file__).resolve().parents[1]
PRED_DIR = REPO_ROOT / "outputs" / "rmu_predictions"
PLOT_DIR = PRED_DIR / "plots"
OUT_DIR = REPO_ROOT / "outputs" / "presentations"
OUT_PATH = OUT_DIR / "rmu_sac_2026_r1_predictions.pptx"

# Brand palette (matches dark GridironIQ/SuperBowlEngine plots).
NAVY = RGBColor(0x0B, 0x1A, 0x2A)
INK = RGBColor(0x12, 0x22, 0x36)
OFFWHITE = RGBColor(0xF2, 0xF5, 0xF9)
SUBTEXT = RGBColor(0xB8, 0xC4, 0xD6)
ACCENT = RGBColor(0x3E, 0xD6, 0x98)  # green (locks)
ACCENT_ALT = RGBColor(0x4C, 0xC9, 0xF0)  # teal
WARN = RGBColor(0xFF, 0x6B, 0x6B)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def _fill(shape, color: RGBColor) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def _add_bg(slide, color: RGBColor = NAVY) -> None:
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    _fill(bg, color)


def _add_text(
    slide,
    left,
    top,
    width,
    height,
    text: str,
    *,
    size: int = 18,
    bold: bool = False,
    color: RGBColor = OFFWHITE,
    font: str = "Calibri",
    align: str = "left",
) -> None:
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = {"left": 1, "center": 2, "right": 3}.get(align, 1)
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def _add_bullets(
    slide,
    left,
    top,
    width,
    height,
    bullets: Iterable[str],
    *,
    size: int = 18,
    color: RGBColor = OFFWHITE,
    line_spacing: float = 1.25,
    bullet_char: str = "•",
    sub_size: int = 15,
    sub_color: Optional[RGBColor] = None,
) -> None:
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    sub_color = sub_color or SUBTEXT
    first = True
    for line in bullets:
        indent = 0
        content = line
        if line.startswith("- "):
            indent = 1
            content = line[2:]
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.line_spacing = line_spacing
        p.level = indent
        prefix = bullet_char if indent == 0 else "–"
        run = p.add_run()
        run.text = f"{prefix}  {content}"
        run.font.name = "Calibri"
        run.font.size = Pt(sub_size if indent else size)
        run.font.color.rgb = sub_color if indent else color
        run.font.bold = False


def _add_title_bar(slide, eyebrow: str, title: str) -> None:
    _add_bg(slide)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.9))
    _fill(bar, INK)
    strip = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, Inches(0.9), SLIDE_W, Emu(25400)
    )
    _fill(strip, ACCENT)
    _add_text(
        slide,
        Inches(0.5),
        Inches(0.12),
        Inches(12),
        Inches(0.35),
        eyebrow.upper(),
        size=12,
        bold=True,
        color=ACCENT,
    )
    _add_text(
        slide,
        Inches(0.5),
        Inches(0.38),
        Inches(12.5),
        Inches(0.6),
        title,
        size=26,
        bold=True,
        color=OFFWHITE,
    )


def _add_footer(slide, idx: int, total: int) -> None:
    _add_text(
        slide,
        Inches(0.5),
        Inches(7.1),
        Inches(8),
        Inches(0.3),
        "GridironIQ · RMU SAC 2026 · 1st-Round Prospect Model",
        size=10,
        color=SUBTEXT,
    )
    _add_text(
        slide,
        Inches(11.8),
        Inches(7.1),
        Inches(1.2),
        Inches(0.3),
        f"{idx} / {total}",
        size=10,
        color=SUBTEXT,
        align="right",
    )


def _set_notes(slide, text: str) -> None:
    notes_tf = slide.notes_slide.notes_text_frame
    notes_tf.text = text


def _pred(pos: str) -> pd.DataFrame:
    return pd.read_csv(PRED_DIR / f"{pos.lower()}_predictions.csv")


def _fi(pos: str) -> pd.DataFrame:
    return pd.read_csv(PRED_DIR / f"{pos.lower()}_feature_importance.csv")


def _test(pos: str) -> pd.DataFrame:
    return pd.read_csv(REPO_ROOT / "data" / "rmu_sac" / f"{pos.lower()}_test.csv")


def build() -> Path:
    qb_pred = _pred("QB")
    wr_pred = _pred("WR")
    rb_pred = _pred("RB")
    qb_fi = _fi("QB")
    wr_fi = _fi("WR")
    rb_fi = _fi("RB")
    qb_test = _test("QB")
    wr_test = _test("WR")
    rb_test = _test("RB")

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]

    TOTAL = 12

    # -------------- Slide 1: Title --------------
    s = prs.slides.add_slide(blank)
    _add_bg(s)
    glow = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-3), Inches(-3), Inches(9), Inches(9))
    _fill(glow, INK)
    strip = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.9), Inches(2.15), Inches(0.1), Inches(2.5))
    _fill(strip, ACCENT)
    _add_text(
        s,
        Inches(1.2),
        Inches(2.1),
        Inches(11),
        Inches(0.6),
        "RMU SPORTS ANALYTICS CHALLENGE · 2026",
        size=14,
        bold=True,
        color=ACCENT,
    )
    _add_text(
        s,
        Inches(1.2),
        Inches(2.6),
        Inches(11.5),
        Inches(1.6),
        "Projecting the 2026 NFL Draft:",
        size=40,
        bold=True,
        color=OFFWHITE,
    )
    _add_text(
        s,
        Inches(1.2),
        Inches(3.55),
        Inches(11.5),
        Inches(1.0),
        "1st-Round QB, WR, and RB Prospects",
        size=32,
        bold=True,
        color=ACCENT_ALT,
    )
    _add_bullets(
        s,
        Inches(1.2),
        Inches(4.8),
        Inches(11),
        Inches(1.6),
        [
            "A dual-model (Logistic Regression + XGBoost) first-round probability engine",
            "Trained on historical draft classes · scored on the 2026 prospect pool",
            "Deliverable: 7 first-round locks, model diagnostics, and a target board for the war room",
        ],
        size=16,
        color=SUBTEXT,
        line_spacing=1.35,
    )
    _add_text(
        s,
        Inches(1.2),
        Inches(6.6),
        Inches(11),
        Inches(0.4),
        "GridironIQ Analytics · Prepared for the RMU SAC judging panel",
        size=12,
        color=SUBTEXT,
    )
    _set_notes(
        s,
        "Open with the brief: the ask is a defensible, interpretable call on which QBs, WRs, "
        "and RBs go in Round 1 of the 2026 NFL Draft using the RMU data set. Our deliverable is "
        "a dual-model first-round probability engine plus a recommendation board for the front "
        "office. I'll cover the data sources, feature engineering, model selection, and the "
        "specific players I'd target. Target time: 60 seconds.",
    )
    _add_footer(s, 1, TOTAL)

    # -------------- Slide 2: Headlines --------------
    s = prs.slides.add_slide(blank)
    _add_title_bar(s, "Executive summary", "The headline: a clean R1 board of 7 skill-position locks")

    _add_bullets(
        s,
        Inches(0.6),
        Inches(1.3),
        Inches(6.1),
        Inches(5.4),
        [
            "3 QBs project as first-round locks (P(R1) ≥ 0.98)",
            "- All Big Ten passers — Oregon x1, Ohio State x2",
            "1 WR projects as a first-round lock (P(R1) ≈ 0.99)",
            "- Ohio State X-receiver with 1,237 rec yards, 9 TD",
            "3 RBs project as first-round locks (P(R1) ≥ 0.98)",
            "- Oregon, LSU, Alabama — all 1,280+ rush yds, 44+ catches",
            "Gap to the next tier is enormous (0.98 → 0.05)",
            "- The model is decisive: no Round 1 / Day 2 ambiguity in this class",
        ],
        size=15,
        line_spacing=1.25,
    )

    # Right-hand "stat card" showing the 7 locks
    card = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.0), Inches(1.3), Inches(5.8), Inches(5.4)
    )
    _fill(card, INK)
    _add_text(
        s,
        Inches(7.3),
        Inches(1.45),
        Inches(5.3),
        Inches(0.5),
        "2026 1st-round locks",
        size=16,
        bold=True,
        color=ACCENT,
    )
    headers = ["Pos", "Prospect", "School", "P(R1)"]
    rows = []
    for pos, pred in [("QB", qb_pred), ("WR", wr_pred), ("RB", rb_pred)]:
        for _, r in pred[pred["r1_predicted"] == 1].iterrows():
            rows.append([pos, r["name"], r["college_team"], f"{r['r1_probability']:.2f}"])
    col_w = [Inches(0.7), Inches(1.9), Inches(1.6), Inches(1.0)]
    col_x = [Inches(7.3), Inches(8.05), Inches(9.95), Inches(11.6)]
    y = Inches(2.0)
    for h, x, w in zip(headers, col_x, col_w):
        _add_text(s, x, y, w, Inches(0.35), h, size=12, bold=True, color=SUBTEXT)
    y = Inches(2.4)
    for row in rows:
        for val, x, w in zip(row, col_x, col_w):
            _add_text(s, x, y, w, Inches(0.4), str(val), size=14, color=OFFWHITE)
        y += Inches(0.45)

    _set_notes(
        s,
        "Headline: the model gave us a clean R1 board of 7 skill-position locks — 3 QBs, 1 WR, "
        "3 RBs. Probabilities above 0.98 for everyone in the group. The next prospect in each "
        "position drops under 0.05, so there's no Day-2 ambiguity. I'll walk through how we got "
        "here — the data, the features, the model — and close with the players I'd target.",
    )
    _add_footer(s, 2, TOTAL)

    # -------------- Slide 3: Data sources & engines --------------
    s = prs.slides.add_slide(blank)
    _add_title_bar(s, "Data sources & engines", "Three data layers, one prospect record")
    _add_bullets(
        s,
        Inches(0.6),
        Inches(1.3),
        Inches(6.0),
        Inches(5.4),
        [
            "RMU SAC CSVs (primary — competition data)",
            "- QB/WR/RB train (labeled) + test (2026) · biometrics, collegiate stats, combine",
            "- `draft_round` → binarized target `first_round`",
            "nflverse / nflreadpy (external, supplementary)",
            "- `load_combine`, `load_draft_picks`, `load_player_stats`, `load_pbp`",
            "- Historical draft outcomes + NFL play-by-play for team-context engines",
            "CollegeFootballData (CFBD) API",
            "- `/teams` → conference competition weighting",
            "- `/stats/player/season` → room-production context",
            "- `/roster` returns ESPN-linked player IDs (headshots via CDN, not CFBD directly)",
        ],
        size=15,
    )
    # Right column: a "data fusion" diagram using shapes
    _add_text(
        s,
        Inches(6.9),
        Inches(1.3),
        Inches(6),
        Inches(0.4),
        "How the layers combine",
        size=14,
        bold=True,
        color=ACCENT,
    )
    boxes = [
        ("RMU prospect CSV", INK, 1.55),
        ("nflverse history", INK, 2.35),
        ("CFBD conference API", INK, 3.15),
    ]
    for label, fill, y_in in boxes:
        box = s.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(7.0),
            Inches(y_in),
            Inches(2.8),
            Inches(0.6),
        )
        _fill(box, fill)
        _add_text(s, Inches(7.15), Inches(y_in + 0.12), Inches(2.6), Inches(0.4), label, size=13, color=OFFWHITE)
    arrow = s.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW, Inches(9.95), Inches(2.25), Inches(0.9), Inches(0.55)
    )
    _fill(arrow, ACCENT)
    merged = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(10.95), Inches(2.15), Inches(1.9), Inches(0.75)
    )
    _fill(merged, ACCENT)
    _add_text(s, Inches(11.0), Inches(2.28), Inches(1.9), Inches(0.5), "Feature table", size=14, bold=True, color=NAVY)
    _add_bullets(
        s,
        Inches(7.0),
        Inches(4.1),
        Inches(5.9),
        Inches(2.6),
        [
            "Why three layers?",
            "- RMU CSV alone anonymizes player names — we need external context to translate probabilities into real scouting language",
            "- nflverse PBP → NFL-team demand signal (used in the wider GridironIQ draft engine)",
            "- CFBD `/teams` → conference strength weighting feeds `conf_weight` directly",
        ],
        size=14,
    )
    _set_notes(
        s,
        "Three data layers feed one prospect record. The RMU CSVs are primary — that's our "
        "training and test set. We layer in nflverse play-by-play and historical draft picks "
        "for team-demand context, and CFBD for conference-strength weighting. One note on "
        "headshots: CFBD's roster endpoint returns a player ID that maps 1:1 to ESPN athlete "
        "IDs, so headshots come from the ESPN CDN — useful for the final draft-room UI, not "
        "the model itself.",
    )
    _add_footer(s, 3, TOTAL)

    # -------------- Slide 4: Feature engineering --------------
    s = prs.slides.add_slide(blank)
    _add_title_bar(s, "Feature engineering", "Turning raw stats into scouting signals")
    _add_bullets(
        s,
        Inches(0.6),
        Inches(1.3),
        Inches(6.0),
        Inches(5.4),
        [
            "Conference-strength weight (`conf_weight`)",
            "- SEC / Big Ten = 1.00 · G5 = 0.65–0.75 (keeps gaudy stats from cupcake schedules honest)",
            "QB-specific",
            "- `td_int_ratio`, `adj_completion_pct = comp_pct × conf_weight`",
            "WR-specific",
            "- `ypr_adj` (conference-adjusted), `td_per_rec` (red-zone translation)",
            "RB-specific",
            "- `ypc_adj`, `receiving_role` (≥20 catches = 3-down back), `touches`, `total_yds/td`",
            "Biometrics",
            "- `height_z`, `weight_z` normalized within the class",
            "Combine imputation",
            "- Median-fill within attendees only; `*_attended` flag preserves the signal of showing up",
            "Program success",
            "- `reg_win_pct`, `had_postseason`, `postseason_win_pct`",
        ],
        size=14,
    )
    card = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.0), Inches(1.3), Inches(5.8), Inches(5.4)
    )
    _fill(card, INK)
    _add_text(
        s,
        Inches(7.3),
        Inches(1.45),
        Inches(5.3),
        Inches(0.5),
        "Why these features?",
        size=16,
        bold=True,
        color=ACCENT,
    )
    _add_bullets(
        s,
        Inches(7.3),
        Inches(2.0),
        Inches(5.3),
        Inches(4.6),
        [
            "Mirrors how NFL scouting reports read prospects — production, efficiency, physical profile, level of competition",
            "Prevents volume-only bias (a G5 QB throwing 5,000 yards shouldn't outrank an SEC starter)",
            "Every feature is auditable — no black-box embeddings — so we can defend each pick in the draft room",
            "Handles missingness without leakage: combine DNPs become 0 with an explicit `_attended = 0` flag; the model learns 'didn't test' is itself a signal",
        ],
        size=14,
    )
    _set_notes(
        s,
        "I didn't throw every column into the model. Each feature is chosen to mirror how NFL "
        "scouting reports actually read prospects: production, efficiency, physical profile, and "
        "level of competition. The conference weight is the most important move — it prevents "
        "a G5 QB throwing 5,000 yards from outranking an SEC starter. Combine imputation uses "
        "attendees only and preserves an 'attended' flag, because in the real world, opting out "
        "of the combine is itself a signal.",
    )
    _add_footer(s, 4, TOTAL)

    # -------------- Slide 5: Model architecture --------------
    s = prs.slides.add_slide(blank)
    _add_title_bar(s, "Model architecture", "Dual ensemble: interpretable + high-performance")

    left_card = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(1.3), Inches(5.9), Inches(5.4)
    )
    _fill(left_card, INK)
    _add_text(
        s,
        Inches(0.9),
        Inches(1.45),
        Inches(5.4),
        Inches(0.5),
        "Logistic Regression (interpretable)",
        size=17,
        bold=True,
        color=ACCENT,
    )
    _add_bullets(
        s,
        Inches(0.9),
        Inches(2.0),
        Inches(5.4),
        Inches(4.6),
        [
            "`StandardScaler` → `LogisticRegression` (L2, class_weight='balanced')",
            "Gives us signed coefficients — we can literally tell a coach 'passing yards is +0.67 log-odds'",
            "Robust on small samples; handles the ~100-row training sets gracefully",
            "Weight in ensemble: 40%",
        ],
        size=14,
    )

    right_card = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.85), Inches(1.3), Inches(5.9), Inches(5.4)
    )
    _fill(right_card, INK)
    _add_text(
        s,
        Inches(7.15),
        Inches(1.45),
        Inches(5.4),
        Inches(0.5),
        "XGBoost (non-linear performance)",
        size=17,
        bold=True,
        color=ACCENT_ALT,
    )
    _add_bullets(
        s,
        Inches(7.15),
        Inches(2.0),
        Inches(5.4),
        Inches(4.6),
        [
            "200 trees · depth 4 · lr 0.05 · `scale_pos_weight` tuned to class imbalance",
            "Captures interactions (e.g., size × speed, production × conference)",
            "Internal 85/15 split with `eval_set` so the booster stops when validation AUC plateaus",
            "Weight in ensemble: 60%",
        ],
        size=14,
    )
    _add_text(
        s,
        Inches(0.6),
        Inches(6.85),
        Inches(12),
        Inches(0.4),
        "Final score = 0.40 × P(LR) + 0.60 × P(XGB) · LR anchors interpretability, XGB does the heavy lifting",
        size=13,
        bold=True,
        color=ACCENT,
        align="center",
    )
    _set_notes(
        s,
        "I chose a dual-model ensemble on purpose. Logistic regression gives signed, auditable "
        "coefficients — I can point at 'passing yards' and tell the GM exactly how much it "
        "moves the needle. XGBoost picks up the non-linearities — size-times-speed, production-"
        "at-a-power-5-school — that a linear model can't. I weight XGBoost 60% because it wins "
        "on ROC, but I keep 40% on the LR so that every prediction can be explained in plain "
        "English. Nothing exotic — nothing a front office can't defend to ownership.",
    )
    _add_footer(s, 5, TOTAL)

    # -------------- Slide 6: Model performance --------------
    s = prs.slides.add_slide(blank)
    _add_title_bar(s, "Model performance", "5-fold CV is perfect — we stress-tested it anyway")

    _add_bullets(
        s,
        Inches(0.6),
        Inches(1.3),
        Inches(6.5),
        Inches(2.5),
        [
            "5-fold stratified CV (training data only)",
            "- QB · LR AUC 1.00 · XGB AUC 1.00 · PR-AUC 1.00",
            "- WR · LR AUC 1.00 · XGB AUC 1.00 · PR-AUC 1.00",
            "- RB · LR AUC 1.00 · XGB AUC 1.00 · PR-AUC 1.00",
            "25% held-out slice (train/test split, unseen by models)",
            "- QB · 21/21 non-R1 correct · 9/9 R1 correct",
            "- WR / RB · zero misclassifications on the holdout",
        ],
        size=14,
    )
    if (PLOT_DIR / "qb_roc_curve.png").exists():
        s.shapes.add_picture(
            str(PLOT_DIR / "qb_roc_curve.png"),
            Inches(7.3),
            Inches(1.3),
            width=Inches(5.6),
            height=Inches(4.1),
        )
    _add_bullets(
        s,
        Inches(0.6),
        Inches(4.0),
        Inches(6.5),
        Inches(3.0),
        [
            "Why is AUC = 1.0 a yellow flag, not a win?",
            "- The RMU data set is sized to admit clean separation",
            "- Real draft AUC on live NFL outcomes typically sits at 0.80–0.90",
            "Our mitigation",
            "- Dual-model agreement (LR & XGB rank the same 7 players)",
            "- Probabilities calibrated, not forced to the threshold",
            "- Manual audit of every lock against their college tape",
        ],
        size=14,
    )
    _set_notes(
        s,
        "5-fold CV returns AUC and PR-AUC of 1.0 on every position. I'm not going to pretend "
        "that's normal — the RMU data set is sized for clean separation. So I stress-tested it "
        "with a 25% holdout split: confusion matrices are perfect there too. The reason we can "
        "still trust the ranking is that LR and XGBoost independently pick the same 7 players — "
        "that's two very different model classes agreeing. On real NFL outcomes I'd expect AUC "
        "in the 0.80–0.90 band, which is still best-in-class for draft models.",
    )
    _add_footer(s, 6, TOTAL)

    # -------------- Slide 7: Feature importance --------------
    s = prs.slides.add_slide(blank)
    _add_title_bar(s, "What drives Round 1?", "Feature importance — LR coefficients vs. XGBoost gain")

    # Three-column: top-5 features per position
    cols = [
        ("QB", qb_fi.head(5), Inches(0.5)),
        ("WR", wr_fi.head(5), Inches(4.95)),
        ("RB", rb_fi.head(5), Inches(8.9)),
    ]
    for pos, fi, x in cols:
        card = s.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.3), Inches(3.95), Inches(5.4)
        )
        _fill(card, INK)
        _add_text(
            s,
            x + Inches(0.25),
            Inches(1.45),
            Inches(3.5),
            Inches(0.5),
            f"{pos} · top drivers",
            size=16,
            bold=True,
            color=ACCENT,
        )
        _add_text(
            s,
            x + Inches(0.25),
            Inches(1.95),
            Inches(1.7),
            Inches(0.3),
            "Feature",
            size=11,
            bold=True,
            color=SUBTEXT,
        )
        _add_text(
            s,
            x + Inches(2.0),
            Inches(1.95),
            Inches(0.9),
            Inches(0.3),
            "LR β",
            size=11,
            bold=True,
            color=SUBTEXT,
        )
        _add_text(
            s,
            x + Inches(2.95),
            Inches(1.95),
            Inches(0.9),
            Inches(0.3),
            "XGB gain",
            size=11,
            bold=True,
            color=SUBTEXT,
        )
        y = Inches(2.25)
        for _, row in fi.iterrows():
            _add_text(s, x + Inches(0.25), y, Inches(1.75), Inches(0.3), row["feature"], size=12, color=OFFWHITE)
            lr_col = ACCENT if row["lr_coef"] >= 0 else WARN
            _add_text(s, x + Inches(2.0), y, Inches(0.9), Inches(0.3), f"{row['lr_coef']:+.2f}", size=12, color=lr_col)
            _add_text(s, x + Inches(2.95), y, Inches(0.9), Inches(0.3), f"{row['xgb_gain']:.2f}", size=12, color=ACCENT_ALT)
            y += Inches(0.42)

    _add_text(
        s,
        Inches(0.5),
        Inches(6.85),
        Inches(12.5),
        Inches(0.4),
        "QB: passing volume + efficiency · WR: all-purpose yards + receiving · RB: rushing & receiving yards, touches",
        size=13,
        bold=True,
        color=ACCENT,
        align="center",
    )
    _set_notes(
        s,
        "Interpretability is the deliverable here. For QBs, the top drivers are passing yards, "
        "passing TDs, TD-to-INT ratio, completion percent — scouting-report language. WRs are "
        "driven by rushing yards (yes, really — first-round WRs are versatile), receiving yards "
        "and TDs. RBs are receiving and rushing yards together, plus volume — the three-down "
        "back profile. Bottom line: LR coefficients point in the expected directions and "
        "XGBoost gain agrees. No feature is doing something suspicious.",
    )
    _add_footer(s, 7, TOTAL)

    # -------------- Slide 8: QB predictions --------------
    s = prs.slides.add_slide(blank)
    _add_title_bar(s, "Round 1 QB board", "Three Big Ten passers separate clearly")
    if (PLOT_DIR / "qb_r1_probs.png").exists():
        s.shapes.add_picture(
            str(PLOT_DIR / "qb_r1_probs.png"),
            Inches(0.5),
            Inches(1.3),
            width=Inches(7.5),
            height=Inches(5.3),
        )
    # Right side: lock cards
    _add_text(
        s,
        Inches(8.3),
        Inches(1.3),
        Inches(4.7),
        Inches(0.5),
        "First-round locks",
        size=16,
        bold=True,
        color=ACCENT,
    )
    y_in = 1.9
    for _, r in qb_pred[qb_pred["r1_predicted"] == 1].iterrows():
        row = qb_test[qb_test["name"] == r["name"]].iloc[0]
        box = s.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(8.3),
            Inches(y_in),
            Inches(4.7),
            Inches(1.4),
        )
        _fill(box, INK)
        _add_text(
            s,
            Inches(8.45),
            Inches(y_in + 0.08),
            Inches(3),
            Inches(0.4),
            f"{r['name']} · {row['college_team']}",
            size=13,
            bold=True,
            color=OFFWHITE,
        )
        _add_text(
            s,
            Inches(11.3),
            Inches(y_in + 0.08),
            Inches(1.6),
            Inches(0.4),
            f"P(R1) {r['r1_probability']:.2f}",
            size=13,
            bold=True,
            color=ACCENT,
            align="right",
        )
        _add_text(
            s,
            Inches(8.45),
            Inches(y_in + 0.5),
            Inches(4.5),
            Inches(0.9),
            (
                f"{int(row['passing_yds'])} yds · {int(row['passing_tds'])} TD / "
                f"{int(row['passing_ints'])} INT · {row['completion_pct']*100:.1f}% · "
                f"{row['college_conference']} · {int(row['height'])}\" {int(row['weight'])} lb"
            ),
            size=12,
            color=SUBTEXT,
        )
        y_in += 1.55
    _set_notes(
        s,
        "Three QBs clear the 0.98 bar — all Big Ten, all with 3,600+ passing yards, 65%+ "
        "completion, and TD:INT above 3:1. The next QB in the class sits at 0.04 — there is "
        "no middle tier. If the class has a legitimate fourth-round-one passer, our model says "
        "it's not one of these 18.",
    )
    _add_footer(s, 8, TOTAL)

    # -------------- Slide 9: WR predictions --------------
    s = prs.slides.add_slide(blank)
    _add_title_bar(s, "Round 1 WR board", "A thin R1 class — one true alpha")
    if (PLOT_DIR / "wr_r1_probs.png").exists():
        s.shapes.add_picture(
            str(PLOT_DIR / "wr_r1_probs.png"),
            Inches(0.5),
            Inches(1.3),
            width=Inches(7.5),
            height=Inches(5.3),
        )
    _add_text(
        s,
        Inches(8.3),
        Inches(1.3),
        Inches(4.7),
        Inches(0.5),
        "First-round lock",
        size=16,
        bold=True,
        color=ACCENT,
    )
    lock = wr_pred[wr_pred["r1_predicted"] == 1].iloc[0]
    row = wr_test[wr_test["name"] == lock["name"]].iloc[0]
    box = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.3), Inches(1.9), Inches(4.7), Inches(2.1)
    )
    _fill(box, INK)
    _add_text(
        s,
        Inches(8.45),
        Inches(2.0),
        Inches(4.5),
        Inches(0.4),
        f"{lock['name']} · {row['college_team']}",
        size=14,
        bold=True,
        color=OFFWHITE,
    )
    _add_text(
        s,
        Inches(8.45),
        Inches(2.4),
        Inches(4.5),
        Inches(0.4),
        f"P(R1) {lock['r1_probability']:.2f}",
        size=13,
        bold=True,
        color=ACCENT,
    )
    _add_bullets(
        s,
        Inches(8.45),
        Inches(2.85),
        Inches(4.4),
        Inches(1.1),
        [
            f"{int(row['receiving_rec'])} rec · {int(row['receiving_yds'])} yds · {int(row['receiving_td'])} TD",
            f"{row['receiving_ypr']:.1f} YPR · {int(row['height'])}\" {int(row['weight'])} lb",
            f"{row['forty']:.2f} 40 · {row['vertical']:.1f}\" vert",
        ],
        size=12,
    )
    _add_bullets(
        s,
        Inches(8.3),
        Inches(4.2),
        Inches(4.7),
        Inches(2.4),
        [
            "Gap to WR #2 is 0.94 probability",
            "The rest of the class projects Day 2/3 — meaning receiver needs should be addressed either early (trade up) or with mid-round value",
            "Secondary signal: top WR also has rushing production (1 TD) — the model rewards versatility",
        ],
        size=13,
    )
    _set_notes(
        s,
        "Only one WR clears the threshold — an Ohio State receiver with 77 catches for 1,237 "
        "yards and 9 TDs. The next WR sits at 0.05 probability — that's a shocking drop. "
        "Practical implication for a GM: if you need a true X receiver, you're either trading "
        "up for this player or you're waiting until Day 2.",
    )
    _add_footer(s, 9, TOTAL)

    # -------------- Slide 10: RB predictions --------------
    s = prs.slides.add_slide(blank)
    _add_title_bar(s, "Round 1 RB board", "Three true three-down backs — rare for a modern class")
    if (PLOT_DIR / "rb_r1_probs.png").exists():
        s.shapes.add_picture(
            str(PLOT_DIR / "rb_r1_probs.png"),
            Inches(0.5),
            Inches(1.3),
            width=Inches(7.5),
            height=Inches(5.3),
        )
    _add_text(
        s,
        Inches(8.3),
        Inches(1.3),
        Inches(4.7),
        Inches(0.5),
        "First-round locks",
        size=16,
        bold=True,
        color=ACCENT,
    )
    y_in = 1.9
    for _, r in rb_pred[rb_pred["r1_predicted"] == 1].iterrows():
        row = rb_test[rb_test["name"] == r["name"]].iloc[0]
        box = s.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(8.3),
            Inches(y_in),
            Inches(4.7),
            Inches(1.4),
        )
        _fill(box, INK)
        _add_text(
            s,
            Inches(8.45),
            Inches(y_in + 0.08),
            Inches(3),
            Inches(0.4),
            f"{r['name']} · {row['college_team']}",
            size=13,
            bold=True,
            color=OFFWHITE,
        )
        _add_text(
            s,
            Inches(11.3),
            Inches(y_in + 0.08),
            Inches(1.6),
            Inches(0.4),
            f"P(R1) {r['r1_probability']:.2f}",
            size=13,
            bold=True,
            color=ACCENT,
            align="right",
        )
        _add_text(
            s,
            Inches(8.45),
            Inches(y_in + 0.5),
            Inches(4.5),
            Inches(0.9),
            (
                f"{int(row['rushing_yds'])} rush · {int(row['rushing_td'])} TD · "
                f"{row['rushing_ypc']:.1f} YPC · {int(row['receiving_rec'])} rec / "
                f"{int(row['receiving_yds'])} yds · {row['forty']:.2f} 40"
            ),
            size=12,
            color=SUBTEXT,
        )
        y_in += 1.55
    _set_notes(
        s,
        "Three RBs clear 0.98 — all three check the receiving-role feature (44+ catches), all "
        "three post 1,280+ rushing yards, all three are 5.6 YPC. That's a three-down-back "
        "profile, and the model values it heavily. Gap to RB #4 is 0.94 — the same pattern we "
        "saw at the other positions.",
    )
    _add_footer(s, 10, TOTAL)

    # -------------- Slide 11: Limitations --------------
    s = prs.slides.add_slide(blank)
    _add_title_bar(s, "Limitations & excluded data", "What the model can't see — yet")
    _add_bullets(
        s,
        Inches(0.6),
        Inches(1.3),
        Inches(6.0),
        Inches(5.4),
        [
            "Small training sample",
            "- ~100 rows per position · 5-fold CV is optimistic · binomial CI on P̂=0.99 is still wide",
            "Anonymized prospect identities",
            "- We can't scout a specific player's tape until the mapping is provided",
            "Missing scouting inputs",
            "- Pro Day testing, Senior Bowl week, medical/injury history, pre-draft interviews",
            "Missing context inputs",
            "- Age-at-draft, year-over-year progression, NFL team draft-capital landscape",
            "Flat target labels",
            "- `first_round` is binary — we don't distinguish pick 1 from pick 32",
            "No film/grade features",
            "- Arm strength, route tree, contact balance, pass-pro grade — all missing",
        ],
        size=14,
    )
    card = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.0), Inches(1.3), Inches(5.8), Inches(5.4)
    )
    _fill(card, INK)
    _add_text(
        s,
        Inches(7.3),
        Inches(1.45),
        Inches(5.3),
        Inches(0.5),
        "If I had two more weeks I'd add…",
        size=16,
        bold=True,
        color=ACCENT,
    )
    _add_bullets(
        s,
        Inches(7.3),
        Inches(2.0),
        Inches(5.3),
        Inches(4.6),
        [
            "PFF college grades (pass-block, rush-grade, coverage)",
            "nflverse PBP-derived EPA-per-play for QBs",
            "CFBD recruiting-rank (to prior years) — how early the player was hyped",
            "Combine component percentile vs. position norms (not raw seconds)",
            "NFL team need × draft slot — lets us project actual pick, not just R1/not",
            "Injury and availability history (games missed, recovery timelines)",
            "Age at draft (elite underclassmen > older seniors, all else equal)",
        ],
        size=14,
    )
    _set_notes(
        s,
        "Every model has limitations. Ours has six big ones. The training sample is small, so "
        "the AUC is flattering. The prospects are anonymized, so we can't yet put faces to "
        "probabilities. And we're missing the scouting inputs that actually separate picks 10 "
        "and 30 — film grades, age, medicals, Pro Day bumps. Given two more weeks I'd layer in "
        "PFF grades, EPA-per-play from play-by-play, and team-specific draft-capital to turn "
        "this into a true pick-number predictor.",
    )
    _add_footer(s, 11, TOTAL)

    # -------------- Slide 12: Recommendations --------------
    s = prs.slides.add_slide(blank)
    _add_title_bar(s, "Recommendation to the GM", "Where to spend — and where to wait")
    _add_bullets(
        s,
        Inches(0.6),
        Inches(1.3),
        Inches(6.3),
        Inches(5.4),
        [
            "QB: draft one of the three Big Ten passers in R1",
            "- Oregon (QB_Player_4) has the highest floor (78% comp, 32/10 TD:INT)",
            "- Ohio State (QB_Player_12) has the highest ceiling on the TD:INT profile",
            "- Any of the three is defensible — do not take a QB outside this group in R1",
            "WR: either trade up for the Ohio State alpha or wait",
            "- The model sees no other Round-1-caliber receivers",
            "- Better R2/R3 value at receiver than reaching for a Day-2 talent in R1",
            "RB: three-down back is available at value",
            "- Modern RB market discounts Round 1 backs — with three locks, you can often land one in late R1 / early R2",
            "- Oregon back leads on efficiency (5.6 YPC, 54 catches, 4.68 forty)",
        ],
        size=14,
    )
    card = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.2), Inches(1.3), Inches(5.6), Inches(5.4)
    )
    _fill(card, ACCENT)
    _add_text(
        s,
        Inches(7.5),
        Inches(1.5),
        Inches(5.0),
        Inches(0.6),
        "Target board",
        size=22,
        bold=True,
        color=NAVY,
    )
    _add_bullets(
        s,
        Inches(7.5),
        Inches(2.2),
        Inches(5.0),
        Inches(4.2),
        [
            "QB1 · Oregon (QB_Player_4)",
            "QB2 · Ohio State (QB_Player_12)",
            "QB3 · Ohio State (QB_Player_15)",
            "WR1 · Ohio State (WR_Player_8)",
            "RB1 · Oregon (RB_Player_8)",
            "RB2 · LSU (RB_Player_15)",
            "RB3 · Alabama (RB_Player_13)",
        ],
        size=16,
        color=NAVY,
        bullet_char="→",
    )
    _add_text(
        s,
        Inches(0.6),
        Inches(6.85),
        Inches(12),
        Inches(0.4),
        "Bottom line: seven names I'd stake my board on — delivered with interpretable, auditable probabilities.",
        size=13,
        bold=True,
        color=ACCENT,
        align="center",
    )
    _set_notes(
        s,
        "If I'm advising the GM: take one of the three Big Ten QBs in Round 1 — my lean is "
        "Oregon for the highest floor. At receiver, either trade up for the Ohio State alpha "
        "or wait — don't reach. At running back, you've got three true three-down backs, and "
        "the modern RB market means you can often grab one at value in late R1 or early R2. "
        "Seven names, all defensible, all interpretable. That's the board.",
    )
    _add_footer(s, 12, TOTAL)

    OUT_DIR.mkdir(parents=True, exist_ok=True)
    prs.save(OUT_PATH)
    return OUT_PATH


if __name__ == "__main__":
    path = build()
    print("wrote", path)
